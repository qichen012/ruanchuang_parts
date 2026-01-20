import os
from PyPDF2 import PdfReader
from pptx import Presentation
from openai import OpenAI

# ================= 配置区域 =================

API_KEY = "sk-941f7d49ac55489e92ed0f1888457b6a"
BASE_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1"
MODEL_NAME = "qwen-vl-max"

# 初始化客户端
client = OpenAI(api_key=API_KEY, base_url=BASE_URL)


# ================= 核心功能函数 =================

def extract_text_from_pdf(file_path):
    """从 PDF 中提取文本"""
    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text += content + "\n"
        print(f"[成功] 已读取 PDF: {file_path}")
    except Exception as e:
        print(f"[错误] 读取 PDF 失败 {file_path}: {e}")
    return text


def extract_text_from_ppt(file_path):
    """从 PPT/PPTX 中提取文本"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            # 遍历每张幻灯片中的形状
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        print(f"[成功] 已读取 PPT: {file_path}")
    except Exception as e:
        print(f"[错误] 读取 PPT 失败 {file_path}: {e}")
    return text


def split_text(text, chunk_size=4000):
    """
    简单的文本分段，防止超过模型上下文限制。
    注意：这里按字符粗略分割，生产环境建议使用 Tokenizer。
    """
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]


def call_summary_api(content):
    """调用模型 API 进行总结"""
    prompt = f"""
    请对以下教学材料或文档内容进行学术风格的总结。
    要求：
    1. 提炼核心概念和关键知识点。
    2. 如果有具体的公式或定义，请重点保留。
    3. 输出格式清晰，每个知识点以这样的格式递归输出为json：
    {{
      "id": "root_01",
      "title": "章节名称",
      "content": "简短的概述",
      "children": [
        {{
          "id": "child_01",
          "title": "子知识点标题",
          "content": "详细内容、公式或定义",
          "relatedNodeId": "如果该点与另一个ID有逻辑关联，请注明其ID，否则填null",
          "children": [] 
        }}
      ]
    }}
                
    内容如下：
    {content}
    """

    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": "你是一个专业的学术助教，擅长总结课程内容。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"[API Error] 调用失败: {e}"


# ================= 主程序逻辑 =================

def process_file(file_path):
    """处理单个文件的主流程"""
    ext = os.path.splitext(file_path)[1].lower()

    # 1. 提取文本
    if ext == '.pdf':
        raw_text = extract_text_from_pdf(file_path)
    elif ext in ['.ppt', '.pptx']:
        raw_text = extract_text_from_ppt(file_path)
    else:
        print(f"[跳过] 不支持的文件格式: {file_path}")
        return

    if not raw_text.strip():
        print("[警告] 未提取到文本，可能是扫描件或图片型文档。")
        return

    # 2. 处理长文本 (Map-Reduce 简化版)
    # 如果文本太长，先分段总结，最后再汇总（这里演示逐段处理）
    chunks = split_text(raw_text)
    print(f"--- 文档已切分为 {len(chunks)} 个片段，开始总结 ---")

    final_summary = []
    for i, chunk in enumerate(chunks):
        print(f"正在处理第 {i + 1}/{len(chunks)} 部分...")
        summary = call_summary_api(chunk)
        final_summary.append(summary)

    # 3. 保存结果
    output_filename = f"summary_{os.path.basename(file_path)}.json"
    with open(output_filename, "w", encoding="utf-8") as f:
        f.write(f"# {os.path.basename(file_path)} 总结报告\n\n")
        f.write("\n\n---\n\n".join(final_summary))

    print(f"[完成] 总结已保存至: {output_filename}")


# 使用示例：直接运行脚本
if __name__ == "__main__":
    # 在这里指定你要处理的文件路径
    target_file = "test.pdf"

    if os.path.exists(target_file):
        process_file(target_file)
    else:
        print(f"文件不存在: {target_file}")