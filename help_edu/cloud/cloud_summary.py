import os
import shutil
from fastapi import FastAPI, UploadFile, File, Form
from pydantic import BaseModel
from PyPDF2 import PdfReader
from pptx import Presentation
from openai import OpenAI
from mem0 import Memory

API_KEY = "sk-941f7d49ac55489e92ed0f1888457b6a"
BASE_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1"
MODEL_NAME = "qwen-vl-max"
client = OpenAI(api_key=API_KEY, base_url=BASE_URL)


mem0_config = {
    "vector_store": {
        "provider": "qdrant",
        "config": {
            "path": "./mem0_storage", # 记忆数据会存在这个文件夹里
        }
    },
    # Mem0 内部也需要调用模型来整理记忆，建议和上面保持一致
    "llm": {
        "provider": "ollama",
        "config": {
            "model": "qwen2.5:7b",
            "temperature": 0.1,
            "max_tokens": 1500,
        }
    },
    "embedder": {
        "provider": "ollama",
        "config": {
            "model": "nomic-embed-text" # 确保你 ollama pull nomic-embed-text
        }
    }
}
memory = Memory.from_config(mem0_config)


app = FastAPI()
# ================= 核心功能函数 =================

def extract_text(file_path):
    """从 PDF/PPT 提取文本"""
    text = ""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif ext in ['.ppt', '.pptx']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        print(f"解析错误: {e}")
    return text


def generate_summary(text):
    """调用大模型生成总结"""
    # 截断一下防止超长 (生产环境建议用 tokenizer)
    input_content = text[:10000]

    prompt = f"""
    请对以下教学课件或文档内容进行详细的结构化总结。
    要求：保留核心定义、公式、关键数据，并使用 Markdown 格式。
    每个知识点以这样的格式递归输出：
    {{
      "id": "root_01",
      "title": "章节名称",
      "content": "简短的概述",
      "children": [
        {{
          "id": "child_01",
          "title": "子知识点标题",
          "content": "详细内容、公式或定义",
          "relatedNodeId": "如果该点与另一个ID有逻辑关联，请注明其ID，否则填null",
          "children": [] 
        }}
      ]
    }}
    内容：
    {input_content}
    """

    try:
        resp = client.chat.completions.create(
            model=LLM_MODEL,
            messages=[
                {"role": "system", "content": "你是一个专业的学术助教。"},
                {"role": "user", "content": prompt}
            ]
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"总结生成失败: {e}"

# ================= 主程序逻辑 =================

@app.post("/upload_and_memorize")
async def upload_document(
        file: UploadFile = File(...),
        user_id: str = Form(...)  # 必须传入 user_id，否则不知道是谁的记忆
):
    """
    1. 接收文件 -> 2. 解析 -> 3. 总结 -> 4. 存入 Mem0
    """
    # 1. 保存临时文件
    temp_filename = f"temp_{file.filename}"
    with open(temp_filename, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # 2. 提取文本
    raw_text = extract_text(temp_filename)

    if len(raw_text) < 50:
        os.remove(temp_filename)
        return {"status": "error", "msg": "文件内容太少或无法解析"}

    # 3. 生成总结
    summary_md = generate_summary(raw_text)

    # 4. 【关键步骤】存入 Mem0
    # metadata 很有用，以后你可以根据 filename 搜索特定的记忆
    print(f"正在将 {file.filename} 的总结存入 Mem0...")
    memory.add(
        summary_md,
        user_id=user_id,
        metadata={"filename": file.filename, "type": "document_summary"}
    )

    # 清理
    os.remove(temp_filename)

    return {
        "status": "success",
        "msg": "文件已学习并存入大脑",
        "summary_preview": summary_md[:200] + "..."
    }


class ChatRequest(BaseModel):
    user_id: str
    question: str


@app.post("/chat")
async def chat_with_brain(req: ChatRequest):
    """
    基于 Mem0 记忆的问答接口
    """
    # 1. Mem0 自动检索相关记忆
    # search 接口会返回 list，里面包含最相关的记忆片段
    relevant_memories = memory.search(req.question, user_id=req.user_id, limit=3)

    # 2. 拼接上下文
    context_str = ""
    for mem in relevant_memories:
        # mem['memory'] 是记忆文本
        context_str += f"- {mem['memory']}\n"

    print(f"检索到的背景知识:\n{context_str}")

    # 3. 如果没找到相关记忆，就普通回答；找到了就基于记忆回答
    if not context_str:
        system_prompt = "你是一个助手。请直接回答用户问题。"
        user_content = req.question
    else:
        system_prompt = "你是一个基于知识库的助教。请根据【记忆片段】回答【用户问题】。"
        user_content = f"【记忆片段】:\n{context_str}\n\n【用户问题】: {req.question}"

    # 4. 调用 LLM 生成最终回复
    resp = client.chat.completions.create(
        model=LLM_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_content}
        ]
    )

    return {"answer": resp.choices[0].message.content, "sources": relevant_memories}